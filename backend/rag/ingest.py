"""Ingestion: extract text, chunk, embed, and write a per-user FAISS index.

Indexes are user-scoped (see ``rag.user_store``); ``user_id=None`` uses the
shared ``_anon`` namespace for the CLI / directory entry point.
"""
from __future__ import annotations

import os
from typing import Dict, List

import numpy as np
from sentence_transformers import SentenceTransformer

try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None  # type: ignore
try:
    from docx import Document
except ImportError:
    Document = None  # type: ignore
try:
    import pandas as pd
except ImportError:
    pd = None  # type: ignore
try:
    from pptx import Presentation
except ImportError:
    Presentation = None  # type: ignore

from config import (
    CHUNK_OVERLAP,
    CHUNK_SIZE,
    DOCS_DIR,
    EMBEDDING_MODEL_NAME,
)
from rag.chunking import recursive_split
from rag.hybrid import build_bm25
from rag.user_store import IndexPaths, paths_for
from rag.vector_store import get_store

model = SentenceTransformer(EMBEDDING_MODEL_NAME)


def embed(texts: List[str]) -> np.ndarray:
    return model.encode(texts, normalize_embeddings=True).astype("float32")


# Document Loading


def extract_text_from_file(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".txt":
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    if ext == ".pdf":
        if PdfReader is None:
            raise ImportError("PyPDF2 is required for PDF files")
        reader = PdfReader(file_path)
        return "\n".join((page.extract_text() or "") for page in reader.pages)

    if ext in (".doc", ".docx"):
        if Document is None:
            raise ImportError("python-docx is required for Word files")
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs)

    if ext == ".csv":
        if pd is None:
            raise ImportError("pandas is required for CSV files")
        return pd.read_csv(file_path).to_string()

    if ext in (".xlsx", ".xls"):
        if pd is None:
            raise ImportError("pandas is required for Excel files")
        return pd.read_excel(file_path).to_string()

    if ext in (".ppt", ".pptx"):
        if Presentation is None:
            raise ImportError("python-pptx is required for PowerPoint files")
        prs = Presentation(file_path)
        out: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    out.append(shape.text)
        return "\n".join(out)

    raise ValueError(f"Unsupported file format: {ext}")


def load_text_files_from_dir(directory: str) -> Dict[str, str]:
    documents: dict[str, str] = {}
    for filename in os.listdir(directory):
        if filename.endswith(".txt"):
            path = os.path.join(directory, filename)
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                documents[filename] = f.read()
    return documents


def load_text_files_from_paths(file_paths: List[str]) -> Dict[str, str]:
    documents: dict[str, str] = {}
    for path in file_paths:
        filename = os.path.basename(path)
        try:
            documents[filename] = extract_text_from_file(path)
        except Exception as e:
            print(f"Error loading {filename}: {e}")
    return documents


# Chunking


def chunk_text(text: str, chunk_size: int, overlap: int) -> List[str]:
    """Recursive character splitter: respects paragraph/sentence boundaries."""
    return recursive_split(text, chunk_size=chunk_size, overlap=overlap)


# Core Ingestion


def _user_id_from_paths(paths: IndexPaths) -> str | None:
    return None if paths.namespace == "_anon" else paths.namespace


def _write_index(paths: IndexPaths, documents: Dict[str, str]) -> None:
    all_chunks: list[str] = []
    metadata: list[dict] = []
    for source, text in documents.items():
        for i, chunk in enumerate(chunk_text(text, CHUNK_SIZE, CHUNK_OVERLAP)):
            all_chunks.append(chunk)
            metadata.append({
                "chunk_id": f"{source}_chunk_{i}",
                "source": source,
                "text": chunk,
            })

    if not all_chunks:
        raise ValueError("No valid text chunks found for ingestion")

    embeddings = embed(all_chunks)

    get_store().upsert(paths.namespace, embeddings, metadata)

    # BM25 still uses local-disk pickle keyed by user_id — works alongside both
    # vector backends and is loaded into RAM on first query.
    from rag.hybrid import reload_bm25
    build_bm25(user_id=_user_id_from_paths(paths), metadata=metadata)
    reload_bm25(_user_id_from_paths(paths))


def ingest_documents(documents: Dict[str, str], *, user_id: str | None = None) -> None:
    _write_index(paths_for(user_id), documents)


# CLI entry (shared _anon namespace)


def ingest() -> None:
    print("Loading documents from DOCS_DIR...")
    docs = load_text_files_from_dir(DOCS_DIR)
    ingest_documents(docs, user_id=None)
    print("Ingestion complete.")


# Backend entry


def ingest_files(file_paths: List[str], *, user_id: str | None = None) -> None:
    print("Loading uploaded files...")
    documents = load_text_files_from_paths(file_paths)
    if not documents:
        raise ValueError("No valid files provided for ingestion")
    ingest_documents(documents, user_id=user_id)
    print("Ingestion from uploaded files complete.")


if __name__ == "__main__":
    ingest()
